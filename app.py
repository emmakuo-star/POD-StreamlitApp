import os 
import io
import streamlit as st
import re
from datetime import datetime
import collections

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

# 頁面基本設定
st.set_page_config(
    page_title="POD Report",
    page_icon="🐍"   
)

# ==========================================
# 1. 輔助與樣式函式
# ==========================================
def set_cell_border(cell, color="000000"):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = OxmlElement(border_name)
        ln.set('w', '12700')  # 寬度 1pt
        ln.set('cmpd', 'sng') # 強制宣告為單一實線 (Solid line)
        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', color)
        solidFill.append(srgbClr)
        ln.append(solidFill)
        tcPr.append(ln)

def translate_to_english(text):
    if not text or text == "NA": return "NA"
    replacements = {
        "評估可行性": "Feasibility Evaluation",
        "PASS (系統能力足以支援需求，且有冗餘)": "PASS",
        "FAIL (流量不足": "FAIL (Insufficient flow ",
        "現有": "The existing ",
        "台": " units ",
        "設備足以支援": " are sufficient to support ",
        "設備，並滿足冗餘條件。": " IT units, meeting redundancy requirements.",
        "設備無法支援": " cannot support ",
        "設備，流量不足。": " IT units due to insufficient flow rate.",
        "為解決LC系統的不足，建議新增": "To resolve LC shortage, recommend adding ",
        "需要增加至少": "Need to add at least ",
        "以滿足風量需求並確保冗餘配置。": "to meet airflow requirements and ensure redundancy.",
        "冷氣足以支援 IT 設備的冷卻需求，並且有冗餘。": "units are sufficient for cooling with redundancy.",
        "冷氣無法滿足 IT 設備的風量需求，差距為": "units cannot meet airflow requirements, gap is "
    }
    for ch, en in replacements.items():
        text = text.replace(ch, en)
    return text.strip()

def add_blank_slide(prs):
    """創建並套用名稱為 'empty' 的模板頁，並強制設為純白背景"""
    empty_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() == 'empty':
            empty_layout = layout
            break
            
    if not empty_layout:
        empty_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        
    slide = prs.slides.add_slide(empty_layout)
    
    for shape in list(slide.shapes): 
        sp = shape._element
        sp.getparent().remove(sp)
        
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    return slide

def style_table(table):
    table.first_row = False
    table.horz_banding = False
    table.first_col = False
    table.last_row = False
    table.last_col = False
    table.vert_banding = False
    
    for row in table.rows:
        for cell in row.cells:
            set_cell_border(cell, "000000") 
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = 'Arial'
                    run.font.size = Pt(9)
            
            if cell.text.strip() == "NA" and cell.text_frame.paragraphs:
                if cell.text_frame.paragraphs[0].runs:
                    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(191, 191, 191)

# ==========================================
# 2. 資料解析器
# ==========================================
def parse_data(raw_text):
    data = collections.defaultdict(lambda: "NA")
    
    data['Client'] = re.search(r'Client:\s*(.+)', raw_text).group(1).strip() if re.search(r'Client:\s*(.+)', raw_text) else "NA"
    data['Location'] = re.search(r'Location:\s*(.+)', raw_text).group(1).strip() if re.search(r'Location:\s*(.+)', raw_text) else "NA"
    data['SF'] = re.search(r'SF:\s*(.+)', raw_text).group(1).strip() if re.search(r'SF:\s*(.+)', raw_text) else "NA"
    data['Redundancy'] = re.search(r'Redundancy:\s*(.+)', raw_text).group(1).strip() if re.search(r'Redundancy:\s*(.+)', raw_text) else "NA"
    
    req_match = re.search(r'【評估需求】\n(.*?)(?=\nSF)', raw_text, re.DOTALL)
    data['eval_req'] = translate_to_english(req_match.group(1).strip()) if req_match else "NA"

    data['IT'] = {}
    for match in re.finditer(r'(IT-\d+)\s+model(?:\s*\(.*?\))?:\s*(.+)', raw_text):
        data['IT'][match.group(1)] = {'model': match.group(2).strip(), 'qty': 'NA', 'ac_kw': 'NA', 'ac_cfm': 'NA', 'lc_kw': 'NA', 'lc_lpm': 'NA'}
    
    for it_k in data['IT'].keys():
        ac_m = re.search(rf'{it_k}\s+AC spec:\s*([0-9.]+)kw,\s*([0-9.]+)cfm', raw_text)
        if ac_m: data['IT'][it_k].update({'ac_kw': ac_m.group(1), 'ac_cfm': ac_m.group(2)})
        lc_m = re.search(rf'{it_k}\s+LC spec:\s*([0-9.]+)kw,\s*([0-9.]+)lpm', raw_text)
        if lc_m: data['IT'][it_k].update({'lc_kw': lc_m.group(1), 'lc_lpm': lc_m.group(2)})
        qty_m = re.search(rf'{it_k}\s+Qty:\s*(\d+)', raw_text)
        if qty_m: data['IT'][it_k]['qty'] = qty_m.group(1)

    data['CRAH'] = []
    crah_blocks = re.split(r'CRAH model:', raw_text)[1:]
    for block in crah_blocks:
        if '【' in block: block = block.split('【')[0]
        c_info = {'model': block.split('\n')[0].strip(), 'qty': 'NA', 'kw': 'NA', 'cfm': 'NA', 'it_eval': 'NA'}
        spec = re.search(r'(?:CRAH spec|CRAH sepc|冷卻能力:\s*[0-9.]+\s*kw\n風量):\s*([0-9.]+)kw,\s*([0-9.]+)cfm|冷卻能力:\s*([0-9.]+)\s*kw\n風量:\s*([0-9.]+)\s*cfm', block)
        if spec: 
            c_info.update({'kw': spec.group(1) or spec.group(3), 'cfm': spec.group(2) or spec.group(4)})
        qty = re.search(r'CRAH Qty:\s*(\d+)', block)
        if qty: c_info['qty'] = qty.group(1)
        it_eval = re.search(r'以IT-\d+ Qty:\s*(\d+)計算AC Gap', block)
        if it_eval: c_info['it_eval'] = it_eval.group(1)
        data['CRAH'].append(c_info)
        
    data['CDU'] = []
    if 'CDU model:' in raw_text:
        cdu_blocks = re.split(r'CDU model:', raw_text)[1:]
        for block in cdu_blocks:
            if '【' in block: block = block.split('【')[0]
            c_info = {'model': block.split('\n')[0].strip(), 'qty': 'NA', 'kw': 'NA', 'lpm': 'NA'}
            spec = re.search(r'CDU spec:\s*([0-9.]+)kw,\s*([0-9.]+)lpm', block)
            if spec: c_info.update({'kw': spec.group(1), 'lpm': spec.group(2)})
            qty = re.search(r'CDU Qty:\s*(\d+)', block)
            if qty: c_info['qty'] = qty.group(1)
            data['CDU'].append(c_info)

    data['AC_res'] = collections.defaultdict(lambda: "NA")
    ac_req = re.search(r'Total AC requirement.*?([0-9.]+)\s*kw,\s*([0-9.]+)\s*cfm|1\.\s*Total requirement[：:]\s*([0-9.]+)\s*kw,\s*([0-9.]+)\s*cfm', raw_text)
    if ac_req: 
        data['AC_res']['req_kw'] = ac_req.group(1) or ac_req.group(3)
        data['AC_res']['req_cfm'] = ac_req.group(2) or ac_req.group(4)
        
    ac_sol = re.search(r'Total AC solution.*?([0-9.]+)\s*kw,\s*([0-9.]+)\s*cfm|2\.\s*Total solution[：:]\s*([0-9.]+)\s*kw,\s*([0-9.]+)\s*cfm', raw_text)
    if ac_sol:
        data['AC_res']['sol_kw'] = ac_sol.group(1) or ac_sol.group(3)
        data['AC_res']['sol_cfm'] = ac_sol.group(2) or ac_sol.group(4)
        
    ac_gap_line = re.search(r'(?:3\.\s*)?AC [Gg]ap[：:]\s*(.*)', raw_text)
    data['AC_res']['gaps'] = []
    if ac_gap_line:
        gaps_raw = ac_gap_line.group(1).split('；')
        for g in gaps_raw:
            match = re.search(r'([+-]?[0-9.]+)\s*kw,\s*([+-]?[0-9.]+)\s*cfm(?:\s*\((.*?)\))?', g)
            if match: data['AC_res']['gaps'].append({'kw': match.group(1), 'cfm': match.group(2), 'model': match.group(3)})
            
    pf = re.search(r'(?:4\.\s*)?可行性(?: \(kw\))?[：:]\s*(PASS|FAIL)', raw_text)
    if pf: data['AC_res']['pass_fail'] = pf.group(1)
    
    conc_ac = re.search(r'5\.\s*結論[：:]\s*(.*?)\n|結論:\s*\n冷卻能力.*?:\s*(.*?)\n風量.*?:\s*(.*?)\n', raw_text)
    if conc_ac: 
        if conc_ac.group(1): data['AC_res']['conclusion'] = translate_to_english(conc_ac.group(1))
        else: data['AC_res']['conclusion'] = translate_to_english(f"{conc_ac.group(2)} {conc_ac.group(3)}")

    data['LC_res'] = collections.defaultdict(lambda: "NA")
    if 'LC 系統' in raw_text:
        lc_req = re.search(r'LC 系統.*?1\.\s*Total requirement[：:]\s*([0-9.]+)\s*kw,\s*([0-9.]+)\s*lpm', raw_text, re.DOTALL)
        if lc_req: data['LC_res'].update({'req_kw': lc_req.group(1), 'req_lpm': lc_req.group(2)})
        lc_sol = re.search(r'LC 系統.*?2\.\s*Total solution[：:]\s*([0-9.]+)\s*kw,\s*([0-9.]+)\s*lpm', raw_text, re.DOTALL)
        if lc_sol: data['LC_res'].update({'sol_kw': lc_sol.group(1), 'sol_lpm': lc_sol.group(2)})
        lc_gap = re.search(r'LC 系統.*?3\.\s*LC Gap[：:]\s*([+-]?[0-9.]+)\s*kw,\s*([+-]?[0-9.]+)\s*lpm', raw_text, re.DOTALL)
        if lc_gap: data['LC_res'].update({'gap_kw': lc_gap.group(1), 'gap_lpm': lc_gap.group(2)})
        lc_pf = re.search(r'LC 系統.*?4\.\s*可行性[：:]\s*(PASS|FAIL)', raw_text, re.DOTALL)
        if lc_pf: data['LC_res']['pass_fail'] = lc_pf.group(1)
        lc_conc = re.search(r'LC 系統.*?5\.\s*結論[：:]\s*(.*?)\n', raw_text, re.DOTALL)
        if lc_conc: data['LC_res']['conclusion'] = translate_to_english(lc_conc.group(1))
    
    sugg = re.search(r'建議:?\n(.*)', raw_text, re.DOTALL)
    data['recommendation'] = translate_to_english(sugg.group(1).strip()) if sugg else "NA"
    return data

# ==========================================
# 3. PPTX 渲染函式
# ==========================================
def generate_pptx(data, template_path, output_path):
    prs = Presentation(template_path)
    
    it_list = list(data['IT'].values())
    it_str = ", ".join([f"{it['model']}*{it['qty']}" for it in it_list])
    crah_str = ", ".join([f"{c['model']}*{c['qty']}" for c in data['CRAH']]) if data['CRAH'] else "NA"
    cdu_str = ", ".join([f"{c['model']}*{c['qty']}" for c in data['CDU']]) if data['CDU'] else "NA"

    slide_1 = prs.slides[0]
    h1_box = slide_1.shapes.add_textbox(Cm(1), Cm(4.5), Cm(23), Cm(3)).text_frame
    h1_box.text = f"{data['Client']} {data['Location']} Run-in Room Evaluation"
    h1_box.paragraphs[0].font.size = Pt(40)
    h1_box.paragraphs[0].font.bold = True
    h1_box.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    note_text = f"Your name\nEKH500\n{datetime.today().strftime('%Y/%m/%d')}"
    note_box = slide_1.shapes.add_textbox(Cm(1), Cm(8), Cm(10), Cm(2)).text_frame
    note_box.text = note_text
    for p in note_box.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(255, 255, 255)

    slide_2 = add_blank_slide(prs)
    
    h2_tf = slide_2.shapes.add_textbox(Cm(0.7), Cm(0.53), Cm(20), Cm(1.5)).text_frame
    h2_tf.text = f"{data['Client']} {data['Location']}, {it_str}"
    h2_tf.paragraphs[0].font.size = Pt(32)
    h2_tf.paragraphs[0].font.bold = True
    
    tf1 = slide_2.shapes.add_textbox(Cm(0.7), Cm(2.0), Cm(20), Cm(2)).text_frame
    tf1.text = f"REQUEST:\n{it_str}, {crah_str}, {cdu_str}\n{data['eval_req']}\nSF: {data['SF']}, Redundancy: {data['Redundancy']}"
    for p in tf1.paragraphs: p.font.size = Pt(11)
    tf1.paragraphs[0].font.bold = True

    table_eval = slide_2.shapes.add_table(3, 6, Cm(0.96), Cm(4.11), Cm(17), Cm(2)).table
    for i, h in enumerate(["", "Total requirement", "Total requirement", "Total solution", "Total solution", "Result"]):
        table_eval.cell(0, i).text = h
        table_eval.cell(0, i).fill.solid()
        table_eval.cell(0, i).fill.fore_color.rgb = RGBColor(0, 0, 0)
        table_eval.cell(0, i).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
    table_eval.cell(1, 0).text, table_eval.cell(2, 0).text = "AC", "LC"
    table_eval.cell(1, 1).text, table_eval.cell(1, 2).text = str(data['AC_res']['req_kw']), str(data['AC_res']['req_cfm'])
    table_eval.cell(1, 3).text, table_eval.cell(1, 4).text = str(data['AC_res']['sol_kw']), str(data['AC_res']['sol_cfm'])
    table_eval.cell(1, 5).text = str(data['AC_res']['pass_fail'])
    
    table_eval.cell(2, 1).text, table_eval.cell(2, 2).text = str(data['LC_res']['req_kw']), str(data['LC_res']['req_lpm'])
    table_eval.cell(2, 3).text, table_eval.cell(2, 4).text = str(data['LC_res']['sol_kw']), str(data['LC_res']['sol_lpm'])
    table_eval.cell(2, 5).text = str(data['LC_res']['pass_fail'])

    for i in range(6):
        table_eval.cell(1, i).fill.solid()
        table_eval.cell(1, i).fill.fore_color.rgb = RGBColor(251, 227, 214) 
        table_eval.cell(2, i).fill.solid()
        table_eval.cell(2, i).fill.fore_color.rgb = RGBColor(220, 234, 247) 
        
    table_eval.cell(1,0).text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 113, 50)
    table_eval.cell(1,0).text_frame.paragraphs[0].font.bold = True
    table_eval.cell(2,0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    table_eval.cell(2,0).text_frame.paragraphs[0].font.bold = True
    
    if data['AC_res']['pass_fail'] == 'FAIL': table_eval.cell(1, 5).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
    elif data['AC_res']['pass_fail'] == 'PASS': table_eval.cell(1, 5).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 176, 80)
    if data['LC_res']['pass_fail'] == 'FAIL': table_eval.cell(2, 5).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
    elif data['LC_res']['pass_fail'] == 'PASS': table_eval.cell(2, 5).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 176, 80)

    tf2 = slide_2.shapes.add_textbox(Cm(0.7), Cm(6.6), Cm(20), Cm(4)).text_frame
    tf2.text = f"CONCLUSION:\n{data['recommendation']}"
    for p in tf2.paragraphs: p.font.size = Pt(11)
    tf2.paragraphs[0].font.bold = True

    cdu_len = len(data['CDU']) if data['CDU'] else 1
    crah_len = len(data['CRAH']) if data['CRAH'] else 1
    spec_rows = 1 + len(it_list) + crah_len + cdu_len
    
    fixed_bottom_cm = 14.54
    table_height_cm = 0.8 * spec_rows
    dynamic_top_cm = fixed_bottom_cm - table_height_cm
    
    table_spec = slide_2.shapes.add_table(spec_rows, 5, Cm(0.7), Cm(dynamic_top_cm), Cm(9.5), Cm(table_height_cm)).table
    table_spec.first_row = False     
    table_spec.horz_banding = False  
    table_spec.first_col = False     
    table_spec.last_row = False
    table_spec.last_col = False
    table_spec.vert_banding = False

    for i, txt in enumerate(["Model", "AC kw", "AC cfm", "LC kw", "LC lpm"]):
        table_spec.cell(0, i).text = txt
        table_spec.cell(0, i).fill.solid()
        if i == 0: table_spec.cell(0, i).fill.fore_color.rgb = RGBColor(0, 0, 0)
        elif i in [1,2]: table_spec.cell(0, i).fill.fore_color.rgb = RGBColor(233, 113, 50)
        else: table_spec.cell(0, i).fill.fore_color.rgb = RGBColor(0, 112, 192)
        table_spec.cell(0, i).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    r_idx = 1
    for it in it_list:
        table_spec.cell(r_idx, 0).text, table_spec.cell(r_idx, 1).text, table_spec.cell(r_idx, 2).text = it['model'], it['ac_kw'], it['ac_cfm']
        table_spec.cell(r_idx, 3).text, table_spec.cell(r_idx, 4).text = it['lc_kw'], it['lc_lpm']
        r_idx += 1
    for cr in data['CRAH']:
        table_spec.cell(r_idx, 0).text, table_spec.cell(r_idx, 1).text, table_spec.cell(r_idx, 2).text = cr['model'], cr['kw'], cr['cfm']
        table_spec.cell(r_idx, 3).text, table_spec.cell(r_idx, 4).text = "NA", "NA"
        r_idx += 1
    if data['CDU']:
        for cdu in data['CDU']:
            table_spec.cell(r_idx, 0).text, table_spec.cell(r_idx, 1).text, table_spec.cell(r_idx, 2).text = cdu['model'], "NA", "NA"
            table_spec.cell(r_idx, 3).text, table_spec.cell(r_idx, 4).text = cdu['kw'], cdu['lpm']
            r_idx += 1
    else:
        table_spec.cell(r_idx, 0).text = table_spec.cell(r_idx, 1).text = table_spec.cell(r_idx, 2).text = table_spec.cell(r_idx, 3).text = table_spec.cell(r_idx, 4).text = "NA"

    style_table(table_eval)
    style_table(table_spec)

    slide_3 = add_blank_slide(prs)
    
    h2_tf3 = slide_3.shapes.add_textbox(Cm(0.7), Cm(0.53), Cm(20), Cm(1.5)).text_frame
    h2_tf3.text = f"{data['Location']}, {it_str}"
    h2_tf3.paragraphs[0].font.size = Pt(32)
    h2_tf3.paragraphs[0].font.bold = True

    tf3 = slide_3.shapes.add_textbox(Cm(0.7), Cm(2.0), Cm(20), Cm(3.5)).text_frame
    tf3.text = f"REQUEST:\n{it_str}, {crah_str}, {cdu_str}\n{data['eval_req']}\nSF: {data['SF']}, Redundancy: {data['Redundancy']}\n\nRESULTS:\nAC: {data['AC_res']['conclusion']}\nLC: {data['LC_res']['conclusion']}"
    for p in tf3.paragraphs: p.font.size = Pt(11)

    y_offset = 5.98
    gaps = data['AC_res']['gaps'] if data['AC_res']['gaps'] else [{'kw': data['AC_res'].get('gap_kw','NA'), 'cfm': data['AC_res'].get('gap_cfm','NA'), 'model': None}]
    
    for gap_idx, gap in enumerate(gaps):
        ac_rows = 4 + len(it_list)
        t_ac = slide_3.shapes.add_table(ac_rows, 4, Cm(1.05), Cm(y_offset), Cm(7.2), Cm(0.8*ac_rows)).table
        t_ac.columns[1].width = Cm(1.08)
        t_ac.first_row = False
        t_ac.horz_banding = False
        
        t_ac.cell(0, 0).text, t_ac.cell(0, 1).text, t_ac.cell(0, 2).text, t_ac.cell(0, 3).text = "AC", "#", "kw", "CFM"
        
        ar = 1
        for it in it_list:
            t_ac.cell(ar, 0).text, t_ac.cell(ar, 1).text = it['model'], str(it['qty'])
            t_ac.cell(ar, 2).text, t_ac.cell(ar, 3).text = str(it['ac_kw']), str(it['ac_cfm'])
            for c in range(4): t_ac.cell(ar, c).fill.solid(); t_ac.cell(ar, c).fill.fore_color.rgb = RGBColor(251, 227, 214)
            t_ac.cell(ar, 1).fill.fore_color.rgb = RGBColor(255, 255, 0)
            if t_ac.cell(ar, 1).text_frame.paragraphs[0].runs:
                t_ac.cell(ar, 1).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
            ar += 1
            
        t_ac.cell(ar, 0).text = "Total AC requirement"
        t_ac.cell(ar, 2).text, t_ac.cell(ar, 3).text = str(data['AC_res']['req_kw']), str(data['AC_res']['req_cfm'])
        for c in range(4): 
            t_ac.cell(ar, c).fill.solid(); t_ac.cell(ar, c).fill.fore_color.rgb = RGBColor(233, 113, 50)
            t_ac.cell(ar, c).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        ar += 1
        
        crah_model = gap['model'] if gap['model'] else (data['CRAH'][0]['model'] if data['CRAH'] else "NA")
        crah_qty = next((c['qty'] for c in data['CRAH'] if c['model'] == crah_model), "NA") if data['CRAH'] else "NA"
        
        t_ac.cell(ar, 0).text, t_ac.cell(ar, 1).text = crah_model, str(crah_qty)
        t_ac.cell(ar, 2).text, t_ac.cell(ar, 3).text = str(data['AC_res']['sol_kw']), str(data['AC_res']['sol_cfm'])
        for c in range(4): 
            t_ac.cell(ar, c).fill.solid(); t_ac.cell(ar, c).fill.fore_color.rgb = RGBColor(233, 113, 50)
            t_ac.cell(ar, c).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        t_ac.cell(ar, 1).fill.fore_color.rgb = RGBColor(255, 255, 0)
        if t_ac.cell(ar, 1).text_frame.paragraphs[0].runs:
                t_ac.cell(ar, 1).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
        ar += 1
        
        t_ac.cell(ar, 0).text = "AC gap"
        t_ac.cell(ar, 2).text, t_ac.cell(ar, 3).text = gap['kw'], gap['cfm']
        for c in range(4): 
            t_ac.cell(ar, c).fill.solid(); t_ac.cell(ar, c).fill.fore_color.rgb = RGBColor(0, 0, 0)
            t_ac.cell(ar, c).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            if "-" in t_ac.cell(ar, c).text: t_ac.cell(ar, c).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        
        style_table(t_ac)
        y_offset += (0.8 * ac_rows + 0.5)

    lc_rows = 4 + len(it_list)
    t_lc = slide_3.shapes.add_table(lc_rows, 4, Cm(8.5), Cm(5.98), Cm(7.2), Cm(0.8*lc_rows)).table
    t_lc.columns[1].width = Cm(1.08)
    t_lc.first_row = False
    t_lc.horz_banding = False
    
    t_lc.cell(0, 0).text, t_lc.cell(0, 1).text, t_lc.cell(0, 2).text, t_lc.cell(0, 3).text = "LC", "#", "kw", "LPM"
    
    lr = 1
    for it in it_list:
        t_lc.cell(lr, 0).text, t_lc.cell(lr, 1).text = it['model'], str(it['qty'])
        t_lc.cell(lr, 2).text, t_lc.cell(lr, 3).text = str(it['lc_kw']), str(it['lc_lpm'])
        for c in range(4): t_lc.cell(lr, c).fill.solid(); t_lc.cell(lr, c).fill.fore_color.rgb = RGBColor(220, 234, 247)
        t_lc.cell(lr, 1).fill.fore_color.rgb = RGBColor(255, 255, 0)
        if t_lc.cell(lr, 1).text_frame.paragraphs[0].runs:
            t_lc.cell(lr, 1).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
        lr += 1
        
    t_lc.cell(lr, 0).text = "Total LC requirement"
    t_lc.cell(lr, 2).text, t_lc.cell(lr, 3).text = str(data['LC_res']['req_kw']), str(data['LC_res']['req_lpm'])
    for c in range(4): 
        t_lc.cell(lr, c).fill.solid(); t_lc.cell(lr, c).fill.fore_color.rgb = RGBColor(0, 112, 192)
        t_lc.cell(lr, c).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    lr += 1
    
    cdu_model = data['CDU'][0]['model'] if data['CDU'] else "NA"
    cdu_qty = data['CDU'][0]['qty'] if data['CDU'] else "NA"
    t_lc.cell(lr, 0).text, t_lc.cell(lr, 1).text = cdu_model, str(cdu_qty)
    t_lc.cell(lr, 2).text, t_lc.cell(lr, 3).text = str(data['LC_res']['sol_kw']), str(data['LC_res']['sol_lpm'])
    for c in range(4): 
        t_lc.cell(lr, c).fill.solid(); t_lc.cell(lr, c).fill.fore_color.rgb = RGBColor(0, 112, 192)
        t_lc.cell(lr, c).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    t_lc.cell(lr, 1).fill.fore_color.rgb = RGBColor(255, 255, 0)
    if t_lc.cell(lr, 1).text_frame.paragraphs[0].runs:
        t_lc.cell(lr, 1).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
    lr += 1
    
    t_lc.cell(lr, 0).text = "LC gap"
    t_lc.cell(lr, 2).text, t_lc.cell(lr, 3).text = str(data['LC_res'].get('gap_kw','NA')), str(data['LC_res'].get('gap_lpm','NA'))
    for c in range(4): 
        t_lc.cell(lr, c).fill.solid(); t_lc.cell(lr, c).fill.fore_color.rgb = RGBColor(0, 0, 0)
        t_lc.cell(lr, c).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        if "-" in t_lc.cell(lr, c).text: t_lc.cell(lr, c).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
    style_table(t_lc)

    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    for slide in slides_list: xml_slides.remove(slide)
    
    new_order = [slides_list[0], slides_list[-2], slides_list[-1], slides_list[1]]
    for slide in new_order: xml_slides.append(slide)

    prs.save(output_path)

# ==========================================
# 4. Streamlit UI 與狀態初始化
# ==========================================
if "input_content" not in st.session_state:
    st.session_state.input_content = ""
if "pptx_data" not in st.session_state:
    st.session_state.pptx_data = None  
if "download_clicked" not in st.session_state:
    st.session_state.download_clicked = False

def on_download_click():
    st.session_state.download_clicked = True

def on_clear_click():
    st.session_state.input_content = ""       
    st.session_state.pptx_data = None         
    st.session_state.download_clicked = False 

st.title("🐍POD Report Generation")

user_input = st.text_area("請貼上 Calc Robot 計算結果:", height=300, key="input_content")
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "Report_sample.pptx")

col1, col2, col3 = st.columns([2, 4, 1])

with col1:
    generate_btn = st.button("生成 PPTX", use_container_width=True)
with col3:
    clear_btn = st.button("清除內容", on_click=on_clear_click, use_container_width=True)

if generate_btn:
    if st.session_state.input_content.strip():
        parsed_data = parse_data(st.session_state.input_content) 
        
        try:
            pptx_buffer = io.BytesIO()
            generate_pptx(parsed_data, TEMPLATE_PATH, pptx_buffer) 
            pptx_buffer.seek(0)
            
            st.session_state.pptx_data = pptx_buffer.getvalue()
            st.session_state.download_clicked = False 

        except Exception as e:
            st.error(f"寫入 PPTX 時發生錯誤：{e}")
    else:
        st.warning("請先貼上資料。")

if st.session_state.pptx_data is not None:
    if st.session_state.download_clicked:
        st.success("✅ 下載成功！")
    else:
        st.success("✅ 檔案成功生成！點擊下方按鈕下載")
        
    st.download_button(
        label="📥 下載檔案",
        data=st.session_state.pptx_data,
        file_name="Updated_Report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        disabled=st.session_state.download_clicked,  
        on_click=on_download_click                   
    )
